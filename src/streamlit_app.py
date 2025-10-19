# ============================================================
# 🔧 Imports
# ============================================================
import pandas as pd
import numpy as np
import streamlit as st
import joblib
from io import BytesIO
import requests
import shap
import matplotlib.pyplot as plt
from scipy.special import expit
from sklearn.base import BaseEstimator, TransformerMixin
from sklearn.preprocessing import StandardScaler, OneHotEncoder
from sklearn.impute import SimpleImputer
from sklearn.compose import ColumnTransformer
from sklearn.pipeline import Pipeline

# Para generar PPTX
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR
except ModuleNotFoundError:
    import subprocess, sys
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR

# ============================================================
# ⚙️ Configuración de la app
# ============================================================
st.set_page_config(layout="wide", page_title="Simulador de contratación de depósitos",
                   initial_sidebar_state="expanded")

# Forzar tema oscuro mediante CSS
st.markdown(
    """
    <style>
    .css-18e3th9 {background-color: #0E1117;}
    .css-1d391kg {color: #FFFFFF;}
    .st-bf {color: #FFFFFF;}
    </style>
    """, unsafe_allow_html=True
)

st.title("Simulador de contratación de depósitos")

# ============================================================
# 🌐 URLs
# ============================================================
CSV_URL = "https://raw.githubusercontent.com/jmiglesias98/JIL-TestDataScience-1/refs/heads/main/data/raw/clientes_20251016.csv"
MODEL_URL = "https://raw.githubusercontent.com/jmiglesias98/JIL-TestDataScience-1/refs/heads/main/models/mejor_modelo_20251016.joblib"

# ============================================================
# 🧹 Clases personalizadas
# ============================================================
class DataCleaner(BaseEstimator, TransformerMixin):
    def __init__(self):
        self.categorical_cols = {
            "job": ["admin.", "unknown", "unemployed", "management", "housemaid",
                    "entrepreneur", "student", "blue-collar", "self-employed",
                    "retired", "technician", "services"],
            "marital": ["married", "divorced", "single"],
            "education": ["unknown", "secondary", "primary", "tertiary"],
            "default": ["yes", "no"],
            "housing": ["yes", "no"],
            "loan": ["yes", "no"],
            "contact": ["unknown", "telephone", "cellular"],
            "month": ["jan", "feb", "mar", "apr", "may", "jun", "jul", "aug", "sep", "oct", "nov", "dec"],
            "poutcome": ["unknown", "other", "failure", "success"]
        }

    def fit(self, X, y=None):
        return self

    def transform(self, X):
        df = X.copy()
        for col, allowed_values in self.categorical_cols.items():
            if col in df.columns:
                mode_value = df[col].mode()[0]
                df[col] = df[col].apply(
                    lambda val: mode_value if pd.isna(val)
                    else ("unknown" if val not in allowed_values and "unknown" in allowed_values
                          else (mode_value if val not in allowed_values else val))
                )
        return df


class PreprocesadorDinamico(BaseEstimator, TransformerMixin):
    def __init__(self, cols_to_drop_after_ohe=None):
        self.cols_to_drop_after_ohe = cols_to_drop_after_ohe
        self.ct = None
        self.feature_names_out_ = None

    def fit(self, X, y=None):
        num_cols = X.select_dtypes(include=["int64", "float64"]).columns.tolist()
        cat_cols = X.select_dtypes(include=["object"]).columns.tolist()

        num_transformer = Pipeline(steps=[
            ("imputer", SimpleImputer(strategy="median")),
            ("scaler", StandardScaler())
        ])

        cat_transformer = Pipeline(steps=[
            ("imputer", SimpleImputer(strategy="most_frequent")),
            ("ohe", OneHotEncoder(drop="first", handle_unknown="ignore", sparse_output=False))
        ])

        self.ct = ColumnTransformer(transformers=[
            ("num", num_transformer, num_cols),
            ("cat", cat_transformer, cat_cols)
        ])

        self.ct.fit(X)
        self.feature_names_out_ = self.ct.get_feature_names_out()

        if self.cols_to_drop_after_ohe:
            self.feature_names_out_ = [
                c for c in self.feature_names_out_ if c not in self.cols_to_drop_after_ohe
            ]
        return self

    def transform(self, X):
        X_t = self.ct.transform(X)
        df = pd.DataFrame(X_t, columns=self.ct.get_feature_names_out())
        if self.cols_to_drop_after_ohe:
            cols_existentes = [c for c in self.cols_to_drop_after_ohe if c in df.columns]
            df = df.drop(columns=cols_existentes, errors="ignore")
        return df.values

    def get_feature_names_out(self):
        """
        Devuelve los nombres de las variables después de la transformación.
        Si el modelo fue cargado desde un joblib antiguo y no tiene
        feature_names_out_ definido, los reconstruye desde el ColumnTransformer.
        """
        # Si ya está definido, lo usamos
        if hasattr(self, "feature_names_out_") and self.feature_names_out_ is not None:
            return self.feature_names_out_

        # Si existe ColumnTransformer, obtenemos directamente los nombres
        if hasattr(self, "ct") and hasattr(self.ct, "get_feature_names_out"):
            try:
                names = self.ct.get_feature_names_out()
                # Limpiar nombres de prefijos de steps (num__ y cat__)
                names = [n.replace("num__", "").replace("cat__", "") for n in names]
                return names
            except Exception:
                pass

        # Si no hay nada, devolvemos una lista vacía (evita crash)
        return []

# ============================================================
# 📥 Funciones de carga
# ============================================================
@st.cache_data
def fetch_url(url):
    r = requests.get(url)
    r.raise_for_status()
    return r.content

@st.cache_data
def load_df_from_bytes(bts):
    return pd.read_csv(BytesIO(bts), sep=";")

@st.cache_data
def load_model_from_bytes(bts):
    return joblib.load(BytesIO(bts))

# ============================================================
# ⚙️ Cargar datos y modelo
# ============================================================
try:
    csv_bytes = fetch_url(CSV_URL)
    df = load_df_from_bytes(csv_bytes)
except Exception as e:
    st.error(f"❌ Error cargando CSV: {e}")
    st.stop()

try:
    model_bytes = fetch_url(MODEL_URL)
    modelo_pipeline = load_model_from_bytes(model_bytes)
except Exception as e:
    st.error(f"❌ Error cargando modelo: {e}")
    st.stop()

features = df.columns.tolist()

# ============================================================
# 🛠️ Controles de configuración
# ============================================================
st.sidebar.header("⚙️ Configuración")

bg_size = st.sidebar.slider(
    "Tamaño muestra background (para explainer)",
    min_value=10, max_value=min(500, len(df)),
    value=min(100, len(df))
)
background = df.sample(bg_size, random_state=42)

row_selector = st.sidebar.selectbox(
    "Selecciona cliente por índice (posición en el CSV)",
    options=list(range(len(df)))
)
base_row = df.iloc[row_selector:row_selector+1].copy()

# ============================================================
# ✏️ Controles What-If
# ============================================================
if "selected_row" not in st.session_state or st.session_state["selected_row"] != row_selector:
    st.session_state["selected_row"] = row_selector
    st.session_state["mod_values"] = {}

col1, col2 = st.columns(2)
new_row = base_row.copy()
numeric_cols = new_row.select_dtypes(include=[np.number]).columns.tolist()
cat_cols = [c for c in features if c not in numeric_cols]

with col1:
    st.subheader("🔢 Numéricos")
    for c in numeric_cols:
        col_min = float(df[c].quantile(0.01))
        col_max = float(df[c].quantile(0.99))
        col_val = float(base_row.iloc[0][c])
        delta = max(abs(col_val)*0.5, 1.0)
        v_min = min(col_min, col_val - delta)
        v_max = max(col_max, col_val + delta)
        step = (v_max-v_min)/100 if v_max>v_min else 1.0
        default_val = st.session_state["mod_values"].get(c, col_val)
        new_val = st.slider(c, min_value=v_min, max_value=v_max, value=default_val, step=step)
        new_row.at[new_row.index[0], c] = new_val
        st.session_state["mod_values"][c] = new_val

with col2:
    st.subheader("🏷️ Categóricas / Otros")
    for c in cat_cols:
        uniques = df[c].dropna().unique().tolist()
        default = base_row.iloc[0][c]
        default_val = st.session_state["mod_values"].get(c, default)
        if len(uniques) <= 20 and len(uniques) > 0:
            new_val = st.selectbox(c, options=uniques, index=uniques.index(default_val) if default_val in uniques else 0)
        else:
            new_val = st.text_input(f"{c} (valor)", value=str(default_val))
        new_row.at[new_row.index[0], c] = new_val
        st.session_state["mod_values"][c] = new_val

# ============================================================
# 🧾 Comparativa
# ============================================================
comparacion = pd.DataFrame({
    "Variable": base_row.columns,
    "Valor original": base_row.iloc[0].values,
    "Valor modificado": new_row.iloc[0].values
})
comparacion = comparacion[comparacion["Variable"] != "campaign"].reset_index(drop=True)

def highlight_changes(row):
    orig = row["Valor original"]
    mod = row["Valor modificado"]
    if pd.isna(orig) or pd.isna(mod):
        return [""]*3
    if isinstance(orig, (int,float)) and isinstance(mod, (int,float)):
        if mod>orig: return ["","background-color: #006666",""]
        elif mod<orig: return ["","background-color: #006666",""]
    else:
        if orig!=mod: return ["","background-color: #006666",""]
    return [""]*3

st.markdown("### 🧾 Comparativa de valores del cliente")
st.dataframe(comparacion.style.apply(highlight_changes, axis=1), use_container_width=True)

# ============================================================
# 🧩 Predicción y SHAP
# ============================================================

cleaner = modelo_pipeline.named_steps["cleaner"]
preprocessor = modelo_pipeline.named_steps["preprocessor"]
model = modelo_pipeline.named_steps["modelo"]

base_row_clean = cleaner.transform(base_row)
new_row_clean = cleaner.transform(new_row)
base_row_preprocessed = preprocessor.transform(base_row_clean)
new_row_preprocessed = preprocessor.transform(new_row_clean)
background_clean = cleaner.transform(background)
background_preprocessed = preprocessor.transform(background_clean)

X_before = preprocessor.transform(base_row).astype(np.float32).reshape(1, -1)
X_after = preprocessor.transform(new_row).astype(np.float32).reshape(1, -1)
background_array = preprocessor.transform(background).astype(np.float32)

with st.spinner("🧠 Calculando valores SHAP..."):
    explainer = shap.Explainer(modelo_pipeline, background_array)
    shap_values_before = explainer(X_before)
    shap_values_after = explainer(X_after)

feat_names = [f.replace("num__", "").replace("cat__", "") for f in preprocessor.get_feature_names_out()]

exp_before = shap.Explanation(
    values=shap_values_before.values[0],
    base_values=np.mean(shap_values_before.base_values),
    data=pd.Series(X_before[0], index=feat_names),
    feature_names=feat_names
)
exp_after = shap.Explanation(
    values=shap_values_after.values[0],
    base_values=np.mean(shap_values_after.base_values),
    data=pd.Series(X_after[0], index=feat_names),
    feature_names=feat_names
)

prob_before = expit(exp_before.base_values + exp_before.values.sum())
prob_after = expit(exp_after.base_values + exp_after.values.sum())

# ============================================================
# 📊 Probabilidades
# ============================================================
st.markdown("### 📊 Probabilidades")
colA, colB = st.columns(2)
colA.markdown(f"<div style='background-color:#1f77b4; padding:10px; border-radius:5px; color:white; text-align:center;'>Probabilidad original<br>{prob_before:.4f}</div>", unsafe_allow_html=True)
colB.markdown(f"<div style='background-color:#ff7f0e; padding:10px; border-radius:5px; color:white; text-align:center;'>Probabilidad modificada<br>{prob_after:.4f}</div>", unsafe_allow_html=True)

# ============================================================
# 💧 Waterfalls con nombres de variables
# ============================================================
st.markdown("### 💧 Waterfall SHAP")
col1, col2 = st.columns(2)

with col1:
    st.subheader("Antes de modificaciones")
    fig1, ax1 = plt.subplots(figsize=(8,6))
    shap.plots.waterfall(exp_before, max_display=10, show=False)
    plt.title("SHAP antes de modificaciones", color="white")
    st.pyplot(fig1)

with col2:
    st.subheader("Después de modificaciones")
    fig2, ax2 = plt.subplots(figsize=(8,6))
    shap.plots.waterfall(exp_after, max_display=10, show=False)
    plt.title("SHAP después de modificaciones", color="white")
    st.pyplot(fig2)

# ============================================================
# 🖨️ PPTX con gráficas idénticas
# ============================================================
def create_pptx_dark_centered(prob_before, prob_after, comparacion_df, fig_before, fig_after):
    import io
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR

    prs = Presentation()

    # Tamaño de diapositiva
    SLIDE_WIDTH = prs.slide_width
    SLIDE_HEIGHT = prs.slide_height

    # 🎨 Colores
    bg_color = RGBColor(240, 240, 240)   # gris claro
    title_bg = RGBColor(0, 102, 102)     # verde cian oscuro
    prob_bg_before = RGBColor(31, 119, 180)  # azul
    prob_bg_after = RGBColor(255, 127, 14)   # naranja
    text_color = RGBColor(255, 255, 255)
    table_fill = RGBColor(230, 230, 230)
    table_text = RGBColor(0, 0, 0)

    # Aplicar fondo gris a todas las diapositivas
    def apply_slide_bg(slide):
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_title_box(slide, text, top):
        width = SLIDE_WIDTH - Inches(1.0)
        left = Inches(0.5)
        height = Inches(0.6)
        box = slide.shapes.add_textbox(left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = title_bg
        tf = box.text_frame
        tf.text = text
        tf.paragraphs[0].font.size = Pt(16)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = text_color
        tf.paragraphs[0].alignment = 1
        return box

    def add_prob_box(slide, left, top, width, height, text, bg):
        box = slide.shapes.add_textbox(left, top, width, height)
        box.fill.solid()
        box.fill.fore_color.rgb = bg
        tf = box.text_frame
        tf.text = text
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in tf.paragraphs:
            p.font.size = Pt(14)
            p.font.color.rgb = RGBColor(255, 255, 255)
            p.alignment = 1
        return box

    def add_table_from_df(slide, df, top, total_width, left_margin_ratio=0.05):
        n_rows, n_cols = df.shape
        left = int(SLIDE_WIDTH * left_margin_ratio)
        width = int(SLIDE_WIDTH * (1 - 2 * left_margin_ratio))

        # 🧮 Calcular altura dinámica según número de filas (máx: hasta borde inferior)
        available_height = SLIDE_HEIGHT - top - Inches(1)
        base_height = Inches(4.5)
        height = min(base_height + (n_rows - 10) * 0.2 * Inches(0.2), available_height) if n_rows > 10 else base_height
        height = max(Inches(2.5), height)  # nunca menor a 2.5"

        table = slide.shapes.add_table(
            rows=n_rows + 1, cols=n_cols,
            left=left, top=top,
            width=width, height=height
        ).table

        # Ajustar tamaño de fuente según número de filas
        font_size = 12
        if n_rows > 12:
            font_size = 10
        if n_rows > 18:
            font_size = 8

        # Cabeceras
        for j, col_name in enumerate(df.columns):
            cell = table.cell(0, j)
            cell.text = str(col_name)
            cell.fill.solid()
            cell.fill.fore_color.rgb = title_bg
            cell.text_frame.paragraphs[0].font.size = Pt(font_size)
            cell.text_frame.paragraphs[0].font.bold = True
            cell.text_frame.paragraphs[0].font.color.rgb = text_color
            cell.text_frame.paragraphs[0].alignment = 1

        # Celdas de datos
        for i in range(n_rows):
            for j in range(n_cols):
                cell = table.cell(i + 1, j)
                cell.text = str(df.iloc[i, j])
                cell.fill.solid()
                cell.fill.fore_color.rgb = table_fill
                cell.text_frame.paragraphs[0].font.size = Pt(font_size)
                cell.text_frame.paragraphs[0].font.color.rgb = table_text
                cell.text_frame.paragraphs[0].alignment = 1

        # Ajustar ancho de columnas proporcionalmente
        col_width = int(width / n_cols)
        for j in range(n_cols):
            table.columns[j].width = col_width

        return table

    def add_figure_slide(prs, fig, title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        apply_slide_bg(slide)
        add_title_box(slide, title_text, top=Inches(0.3))
        img_stream = io.BytesIO()
        fig.savefig(img_stream, bbox_inches="tight", facecolor="white")
        img_stream.seek(0)
        slide.shapes.add_picture(img_stream, Inches(0.5), Inches(1.0),
                                 width=SLIDE_WIDTH - Inches(1.0),
                                 height=SLIDE_HEIGHT - Inches(1.5))
        return slide

    # --- Slide 1 ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    apply_slide_bg(slide1)
    add_title_box(slide1, "Resultados de la simulación", top=Inches(0.3))
    margin = Inches(0.5)
    box_width = (SLIDE_WIDTH - 3 * margin) / 2
    box_height = Inches(0.8)
    add_prob_box(slide1, margin, Inches(1.0), box_width, box_height,
                 f"Probabilidad original: {prob_before:.4f}", prob_bg_before)
    add_prob_box(slide1, 2 * margin + box_width, Inches(1.0), box_width, box_height,
                 f"Probabilidad modificada: {prob_after:.4f}", prob_bg_after)
    add_table_from_df(slide1, comparacion_df, top=Inches(2.0), total_width=SLIDE_WIDTH)

    # --- Slide 2 y 3 ---
    add_figure_slide(prs, fig_before, "Valores SHAP antes de modificaciones")
    add_figure_slide(prs, fig_after, "Valores SHAP después de modificaciones")

    # Exportar
    pptx_stream = io.BytesIO()
    prs.save(pptx_stream)
    pptx_stream.seek(0)
    return pptx_stream


file_name = f"simulacion_resultados_cliente_{row_selector}.pptx"
pptx_stream = create_pptx_dark_centered(prob_before, prob_after, comparacion, fig1, fig2)

st.download_button(
    label="📥 Descargar  PPTX",
    data=pptx_stream,
    file_name=file_name,
    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
)

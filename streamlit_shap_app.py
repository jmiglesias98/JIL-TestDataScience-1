# ============================================================
# 🔧 Imports
# ============================================================
import pandas as pd
import numpy as np
import streamlit as st
import joblib
from io import BytesIO
import requests
import shap
import matplotlib.pyplot as plt
from scipy.special import expit
from sklearn.base import BaseEstimator, TransformerMixin
from sklearn.preprocessing import StandardScaler, OneHotEncoder
from sklearn.impute import SimpleImputer
from sklearn.compose import ColumnTransformer
from sklearn.pipeline import Pipeline
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import tempfile
import os

# ============================================================
# ⚙️ Configuración de la app
# ============================================================
st.set_page_config(layout="wide", page_title="Simulador de contratación de depósitos")

# Forzar modo oscuro
st.markdown("""
    <style>
        [data-testid="stAppViewContainer"] {
            background-color: #0e1117;
            color: #fafafa;
        }
        [data-testid="stSidebar"] {
            background-color: #111418;
        }
        h1, h2, h3, h4 {
            color: #f8f9fa;
        }
        .section-box {
            background-color: #1f2937;
            padding: 0.6em 1em;
            border-radius: 0.5em;
            margin-top: 1.2em;
            margin-bottom: 0.8em;
            font-weight: 600;
            color: #f8fafc;
        }
        .prob-box {
            padding: 1em;
            border-radius: 0.5em;
            font-size: 1.2em;
            text-align: center;
            font-weight: 600;
        }
        .prob-before {
            background-color: #1e3a8a;
            color: #e0e7ff;
        }
        .prob-after {
            background-color: #065f46;
            color: #d1fae5;
        }
    </style>
""", unsafe_allow_html=True)

st.title("💼 Simulador de contratación de depósitos")

# ============================================================
# 🌐 URLs
# ============================================================
CSV_URL = "https://raw.githubusercontent.com/jmiglesias98/DataScience/refs/heads/main/clientes_20251016.csv"
MODEL_URL = "https://raw.githubusercontent.com/jmiglesias98/DataScience/refs/heads/main/mejor_modelo_20251016.joblib"

# ============================================================
# 🧹 Clases personalizadas
# ============================================================
class DataCleaner(BaseEstimator, TransformerMixin):
    def __init__(self):
        self.categorical_cols = {
            "job": ["admin.", "unknown", "unemployed", "management", "housemaid",
                    "entrepreneur", "student", "blue-collar", "self-employed",
                    "retired", "technician", "services"],
            "marital": ["married", "divorced", "single"],
            "education": ["unknown", "secondary", "primary", "tertiary"],
            "default": ["yes", "no"],
            "housing": ["yes", "no"],
            "loan": ["yes", "no"],
            "contact": ["unknown", "telephone", "cellular"],
            "month": ["jan", "feb", "mar", "apr", "may", "jun", "jul", "aug", "sep", "oct", "nov", "dec"],
            "poutcome": ["unknown", "other", "failure", "success"]
        }

    def fit(self, X, y=None):
        return self

    def transform(self, X):
        df = X.copy()
        for col, allowed_values in self.categorical_cols.items():
            if col in df.columns:
                mode_value = df[col].mode()[0]
                df[col] = df[col].apply(
                    lambda val: mode_value if pd.isna(val)
                    else ("unknown" if val not in allowed_values and "unknown" in allowed_values
                          else (mode_value if val not in allowed_values else val))
                )
        return df


class PreprocesadorDinamico(BaseEstimator, TransformerMixin):
    def __init__(self, cols_to_drop_after_ohe=None):
        self.cols_to_drop_after_ohe = cols_to_drop_after_ohe
        self.ct = None
        self.feature_names_out_ = None

    def fit(self, X, y=None):
        num_cols = X.select_dtypes(include=["int64", "float64"]).columns.tolist()
        cat_cols = X.select_dtypes(include=["object"]).columns.tolist()

        num_transformer = Pipeline(steps=[
            ("imputer", SimpleImputer(strategy="median")),
            ("scaler", StandardScaler())
        ])

        cat_transformer = Pipeline(steps=[
            ("imputer", SimpleImputer(strategy="most_frequent")),
            ("ohe", OneHotEncoder(drop="first", handle_unknown="ignore", sparse_output=False))
        ])

        self.ct = ColumnTransformer(transformers=[
            ("num", num_transformer, num_cols),
            ("cat", cat_transformer, cat_cols)
        ])

        self.ct.fit(X)
        self.feature_names_out_ = self.ct.get_feature_names_out()

        if self.cols_to_drop_after_ohe:
            self.feature_names_out_ = [
                c for c in self.feature_names_out_ if c not in self.cols_to_drop_after_ohe
            ]
        return self

    def transform(self, X):
        X_t = self.ct.transform(X)
        df = pd.DataFrame(X_t, columns=self.ct.get_feature_names_out())
        if self.cols_to_drop_after_ohe:
            cols_existentes = [c for c in self.cols_to_drop_after_ohe if c in df.columns]
            df = df.drop(columns=cols_existentes, errors="ignore")
        return df.values

    def get_feature_names_out(self):
        return self.feature_names_out_

# ============================================================
# 📥 Funciones de carga
# ============================================================
@st.cache_data
def fetch_url(url):
    r = requests.get(url)
    r.raise_for_status()
    return r.content

@st.cache_data
def load_df_from_bytes(bts):
    return pd.read_csv(BytesIO(bts), sep=";")

@st.cache_data
def load_model_from_bytes(bts):
    return joblib.load(BytesIO(bts))

# ============================================================
# ⚙️ Cargar datos y modelo
# ============================================================
try:
    csv_bytes = fetch_url(CSV_URL)
    df = load_df_from_bytes(csv_bytes)
except Exception as e:
    st.error(f"❌ Error cargando CSV: {e}")
    st.stop()

try:
    model_bytes = fetch_url(MODEL_URL)
    modelo_pipeline = load_model_from_bytes(model_bytes)
except Exception as e:
    st.error(f"❌ Error cargando modelo: {e}")
    st.stop()

features = df.columns.tolist()

# ============================================================
# 🛠️ Configuración lateral
# ============================================================
st.sidebar.header("⚙️ Configuración")

bg_size = st.sidebar.slider(
    "Tamaño muestra background (para explainer)",
    min_value=10, max_value=min(500, len(df)),
    value=min(100, len(df))
)
background = df.sample(bg_size, random_state=42)

row_selector = st.sidebar.selectbox(
    "Selecciona cliente por índice (posición en el CSV)",
    options=list(range(len(df)))
)
base_row = df.iloc[row_selector:row_selector+1].copy()

# ============================================================
# ✏️ Controles What-If (con reinicio al cambiar cliente)
# ============================================================
if "selected_row" not in st.session_state or st.session_state["selected_row"] != row_selector:
    st.session_state["selected_row"] = row_selector
    st.session_state["mod_values"] = {}

st.markdown('<div class="section-box">✏️ Controles What-If</div>', unsafe_allow_html=True)

col1, col2 = st.columns(2)
new_row = base_row.copy()
numeric_cols = new_row.select_dtypes(include=[np.number]).columns.tolist()
cat_cols = [c for c in features if c not in numeric_cols]

with col1:
    for c in numeric_cols:
        col_min = float(df[c].quantile(0.01))
        col_max = float(df[c].quantile(0.99))
        col_val = float(base_row.iloc[0][c])
        delta = max(abs(col_val)*0.5, 1.0)
        v_min = min(col_min, col_val - delta)
        v_max = max(col_max, col_val + delta)
        step = (v_max-v_min)/100 if v_max>v_min else 1.0
        default_val = st.session_state["mod_values"].get(c, col_val)
        new_val = st.slider(c, min_value=v_min, max_value=v_max, value=default_val, step=step)
        new_row.at[new_row.index[0], c] = new_val
        st.session_state["mod_values"][c] = new_val

with col2:
    for c in cat_cols:
        uniques = df[c].dropna().unique().tolist()
        default = base_row.iloc[0][c]
        if len(uniques) <= 20:
            default_val = st.session_state["mod_values"].get(c, default)
            new_val = st.selectbox(c, options=uniques, index=uniques.index(default_val) if default_val in uniques else 0)
        else:
            new_val = st.text_input(f"{c} (valor)", value=str(default))
        new_row.at[new_row.index[0], c] = new_val
        st.session_state["mod_values"][c] = new_val

# ============================================================
# 🧾 Tabla comparativa con resaltado de cambios (sin 'campaign')
# ============================================================
st.markdown('<div class="section-box">🧾 Comparativa de valores del cliente</div>', unsafe_allow_html=True)

comparacion = pd.DataFrame({
    "Variable": base_row.columns,
    "Valor original": base_row.iloc[0].values,
    "Valor modificado": new_row.iloc[0].values
})
comparacion = comparacion[comparacion["Variable"] != "campaign"].reset_index(drop=True)

def highlight_changes(row):
    orig, mod = row["Valor original"], row["Valor modificado"]
    if pd.isna(orig) or pd.isna(mod): return [""]*3
    if isinstance(orig, (int, float)) and isinstance(mod, (int, float)):
        if mod > orig: return ["", "", "background-color: #14532d"]
        elif mod < orig: return ["", "", "background-color: #7f1d1d"]
    elif orig != mod:
        return ["", "", "background-color: #78350f"]
    return [""]*3

st.dataframe(comparacion.style.apply(highlight_changes, axis=1), use_container_width=True)

# ============================================================
# 🔍 Predicción + SHAP
# ============================================================
cleaner = modelo_pipeline.named_steps["cleaner"]
preprocessor = modelo_pipeline.named_steps["preprocessor"]
model = modelo_pipeline.named_steps[list(modelo_pipeline.named_steps.keys())[-1]]

base_row_clean = cleaner.transform(base_row)
new_row_clean = cleaner.transform(new_row)

base_row_preprocessed = preprocessor.transform(base_row_clean)
new_row_preprocessed = preprocessor.transform(new_row_clean)

background_array = preprocessor.transform(cleaner.transform(df.sample(bg_size, random_state=42)))

explainer = shap.Explainer(model, background_array)
shap_before = explainer(base_row_preprocessed)
shap_after = explainer(new_row_preprocessed)

feat_names = preprocessor.get_feature_names_out()

exp_before = shap.Explanation(
    values=shap_before.values[0],
    base_values=explainer.expected_value,
    data=base_row_preprocessed[0],
    feature_names=feat_names
)
exp_after = shap.Explanation(
    values=shap_after.values[0],
    base_values=explainer.expected_value,
    data=new_row_preprocessed[0],
    feature_names=feat_names
)

prob_before = expit(exp_before.base_values + exp_before.values.sum())
prob_after = expit(exp_after.base_values + exp_after.values.sum())

# ============================================================
# 📊 Mostrar probabilidades con colores
# ============================================================
st.markdown('<div class="section-box">📈 Comparativa de Probabilidades</div>', unsafe_allow_html=True)
colA, colB = st.columns(2)
colA.markdown(f'<div class="prob-box prob-before">Probabilidad original: {prob_before:.4f}</div>', unsafe_allow_html=True)
colB.markdown(f'<div class="prob-box prob-after">Probabilidad modificada: {prob_after:.4f} ({(prob_after-prob_before)*100:+.2f} pp)</div>', unsafe_allow_html=True)

# ============================================================
# 💧 Waterfall plots
# ============================================================
st.markdown('<div class="section-box">💧 Impacto de variables (SHAP)</div>', unsafe_allow_html=True)
col1, col2 = st.columns(2)
with col1:
    fig1, ax1 = plt.subplots(figsize=(8, 6))
    shap.plots.waterfall(exp_before, max_display=10, show=False)
    st.pyplot(fig1)
with col2:
    fig2, ax2 = plt.subplots(figsize=(8, 6))
    shap.plots.waterfall(exp_after, max_display=10, show=False)
    st.pyplot(fig2)

# ============================================================
# 📤 Exportar PowerPoint
# ============================================================
st.markdown('<div class="section-box">📤 Descargar resultados</div>', unsafe_allow_html=True)

if st.button("💾 Descargar presentación PowerPoint"):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    # Título
    title = slide.shapes.title
    slide.shapes.title.text = "Resultados de la simulación"
    title.text_frame.paragraphs[0].font.bold = True

    # Probabilidades
    left, top = Inches(0.5), Inches(1.5)
    txBox = slide.shapes.add_textbox(left, top, Inches(9), Inches(1))
    tf = txBox.text_frame
    tf.text = f"Probabilidad original: {prob_before:.4f}\nProbabilidad modificada: {prob_after:.4f}"

    # Guardar gráficos temporales
    tmpdir = tempfile.mkdtemp()
    fig1_path = os.path.join(tmpdir, "before.png")
    fig2_path = os.path.join(tmpdir, "after.png")
    fig1.savefig(fig1_path, bbox_inches="tight")
    fig2.savefig(fig2_path, bbox_inches="tight")

    slide2 = prs.slides.add_slide(prs.slide_layouts[5])
    slide2.shapes.title.text = "Impacto de variables (SHAP)"
    slide2.shapes.add_picture(fig1_path, Inches(0.5), Inches(1.5), height=Inches(3))
    slide2.shapes.add_picture(fig2_path, Inches(5), Inches(1.5), height=Inches(3))

    # Resumen
    slide3 = prs.slides.add_slide(prs.slide_layouts[5])
    slide3.shapes.title.text = "Resumen de cambios"
    table_data = comparacion.head(10)
    left, top, width, height = Inches(0.5), Inches(1.5), Inches(9), Inches(4)
    rows, cols = table_data.shape
    table = slide3.shapes.add_table(rows+1, cols, left, top, width, height).table
    for i, col in enumerate(table_data.columns):
        table.cell(0, i).text = col
    for r in range(rows):
        for c in range(cols):
            table.cell(r+1, c).text = str(table_data.iloc[r, c])

    pptx_path = os.path.join(tmpdir, "simulacion_resultados.pptx")
    prs.save(pptx_path)

    with open(pptx_path, "rb") as f:
        st.download_button(
            label="⬇️ Descargar PowerPoint",
            data=f,
            file_name="simulacion_resultados.pptx",
            mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
